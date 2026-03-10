import os
import json
import logging
from datetime import datetime
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import chardet
from PyPDF2 import PdfReader
import PyPDF2
import docx
from fpdf import FPDF

from crewai import Agent, Task, LLM
from typing import Dict, List, Optional
from dotenv import load_dotenv

from langchain.vectorstores import Chroma
from langchain.embeddings import HuggingFaceEmbeddings
from langchain.docstore.document import Document
from langchain.document_loaders import PyPDFLoader, UnstructuredWordDocumentLoader, TextLoader
from langchain.text_splitter import CharacterTextSplitter

load_dotenv()

# Konfigurasi logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

# Ambil GEMINI_API_KEY dari variabel lingkungan
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
if GEMINI_API_KEY is None:
    logging.error("GEMINI_API_KEY tidak ditemukan di variabel lingkungan.")
    raise ValueError("GEMINI_API_KEY harus di-set pada variabel lingkungan.")

# Database global untuk agen dan percakapan
agents_db: Dict[str, Agent] = {}
conversation_db: Dict[str, List[Dict]] = {}
sessions: Dict[str, "OrchestratorAgent"] = {}

AVAILABLE_MODELS = [
    "gpt-3.5-turbo",
    "gpt-4",
    "gemini/gemini-1.0-pro",
    "gemini/gemini-1.5-pro",
    "gemini/gemini-1.5-flash",
    "claude-3-opus",
    "claude-3-sonnet"
]

PROJECT_ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__)))

# Membuat folder untuk menyimpan dokumen
UPLOAD_FOLDER = os.path.join(PROJECT_ROOT, "uploads")
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

OUTPUT_FOLDER = os.path.join(PROJECT_ROOT, "outputs")
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

# Konfigurasi RAG
PERSIST_DIRECTORY = os.path.join(PROJECT_ROOT, 'db')
os.makedirs(PERSIST_DIRECTORY, exist_ok=True)

# Inisialisasi embedding model
embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2", cache_folder='./my_cache')

# Inisialisasi Chroma DB
vectorstore = Chroma(persist_directory=PERSIST_DIRECTORY, embedding_function=embedding_model)

# Kelas RAG untuk mengelola dokumen dan retrieval
class RAGManager:
    def __init__(self):
        self.text_splitter = CharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        
    def add_document(self, content, metadata=None):
        """Menambahkan dokumen ke vectorstore"""
        if metadata is None:
            metadata = {"source": "uploaded_document"}
            
        # Split dokumen menjadi chunks
        texts = self.text_splitter.split_text(content)
        
        try:
            # Buat dokumen langchain
            documents = [Document(page_content=t, metadata=metadata) for t in texts]
            
            # Tambahkan ke vectorstore
            vectorstore.add_documents(documents)
            vectorstore.persist()
            
            logging.info(f"Berhasil menambahkan {len(documents)} chunks ke vectorstore")
            return len(documents)
        
        except Exception as e:
            logging.error(f"Error saat analisis kueri: {e}")
            return "Belum ada dokumen"
    
    def query_documents(self, query, k=5):
        """Mengambil dokumen yang relevan berdasarkan query"""
        docs = vectorstore.similarity_search(query, k=k)
        return docs
    
    def get_relevant_context(self, query, k=5):
        """Mendapatkan konteks yang relevan untuk RAG"""
        docs = self.query_documents(query, k)
        context = "\n\n".join([doc.page_content for doc in docs])
        return context

# Inisialisasi RAG Manager
rag_manager = RAGManager()

# Fungsi untuk mengekstrak teks dari berbagai format dokumen
def extract_text_from_document(file_path):
    """Ekstrak teks dari berbagai format dokumen"""
    file_extension = os.path.splitext(file_path.lower())[1]
    extracted_text = ""
    
    try:
        # Untuk PDF
        if file_extension == '.pdf':
            reader = PdfReader(file_path)
            for page in reader.pages:
                extracted_text += page.extract_text() + "\n"
        
        # Untuk DOCX
        elif file_extension == '.docx':
            from docx import Document
            doc = Document(file_path)
            for para in doc.paragraphs:
                extracted_text += para.text + "\n"
        
        # Untuk PowerPoint (PPT/PPTX)
        elif file_extension in ['.ppt', '.pptx']:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted_text += shape.text + "\n"
                extracted_text += "\n---\n"  # Pembatas antar slide
        
        # Fallback untuk file txt atau lainnya
        else:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                extracted_text = f.read()
    
    except Exception as e:
        logging.error(f"Gagal mengekstrak teks: {e}")
        raise
        
    return extracted_text

def save_to_word(text: str, filename="output.docx"):
    doc = Document()
    doc.add_paragraph(text)
    doc.save(filename)

    return filename

def save_to_pdf(text: str, filename="output.pdf"):
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, text)
    pdf.output(filename)
    
    return filename

class ConversationMemory:
    def __init__(self, session_id: str = "default"):
        self.session_id = session_id
        if session_id not in conversation_db:
            conversation_db[session_id] = []

    def add_interaction(self, user_query: str, agent_response: str, agent_id: str = "orchestrator") -> None:
        conversation_db[self.session_id].append({
            "timestamp": datetime.now().isoformat(),
            "user_query": user_query,
            "agent_response": agent_response,
            "agent_id": agent_id
        })

    def get_recent_history(self, max_turns: int = 5) -> str:
        interactions = conversation_db.get(self.session_id, [])
        recent = interactions[-max_turns:] if interactions else []
        return "\n".join([f"User: {item['user_query']}\nAgent: {item['agent_response']}" for item in recent])

    def save_to_file(self) -> None:
        os.makedirs("conversations", exist_ok=True)
        filename = os.path.join("conversations", f"{self.session_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json")
        with open(filename, "w") as f:
            json.dump(conversation_db[self.session_id], f, indent=2)
        logging.info(f"Percakapan berhasil disimpan ke {filename}")

class OrchestratorAgent:
    def __init__(self, model: str = "gemini/gemini-1.5-flash", session_id: str = "default"):
        self.model = model
        self.llm = LLM(model=model, api_key=GEMINI_API_KEY)
        self.session_id = session_id
        self.memory = ConversationMemory(session_id)
        self.agent = Agent(
            role="Orchestrator",
            goal="Mengatur routing kueri dan mengagregasi respons dari agen lain sambil menjaga konteks percakapan",
            backstory="Anda adalah AI orkestrator yang cerdas, bertugas mengarahkan kueri, mendelegasikan ke agen lain, dan memberikan jawaban yang kontekstual.",
            verbose=True,
            allow_delegation=True,
            llm=self.llm
        )
        agents_db["orchestrator"] = self.agent

    def analyze_query(self, query: str) -> dict:
        history = self.memory.get_recent_history(3)
        # Add clearer instructions to ensure consistent output format:
        task_description = (
            f"Analisis kueri berikut dalam konteks percakapan:\n\n"
            f"Percakapan terakhir:\n{history}\n\n"
            f"Kueri: \"{query}\"\n\n"
            "Klasifikasikan kueri sebagai 'general', 'follow_up', 'needs_agent', atau 'needs_retrieval'. "
            "Jika pertanyaan mengacu pada informasi dalam dokumen, klasifikasikan sebagai 'needs_retrieval'. "
            "Kembalikan output berupa JSON dengan kunci: query_type, agent_needed, reasoning, "
            "requires_context, refers_to_previous, dan needs_retrieval."
            "PENTING: Untuk kunci 'agent_needed', kembalikan string nama domain seperti 'programming', 'science', dll., bukan boolean."
            "Gunakan format yang tepat tanpa Markdown atau simbol seperti '**' untuk bold."
        )
        analysis_task = Task(
            description=task_description,
            agent=self.agent,
            expected_output="JSON analysis of query type with context awareness"
        )
        try:
            result = self.agent.execute_task(analysis_task)
            json_start = result.find("{")
            json_end = result.rfind("}") + 1
            if json_start == -1 or json_end == -1:
                raise ValueError("Output JSON tidak ditemukan.")
            analysis = json.loads(result[json_start:json_end])
            analysis.setdefault("query_type", "general")
            analysis.setdefault("agent_needed", "none")
            analysis.setdefault("requires_context", False)
            analysis.setdefault("refers_to_previous", False)
            analysis.setdefault("needs_retrieval", False)
            return analysis
        except Exception as e:
            logging.error(f"Error saat analisis kueri: {e}")
            return {
                "query_type": "general",
                "agent_needed": "none",
                "reasoning": f"Fallback karena error: {e}",
                "requires_context": any(kw in query.lower() for kw in ["previous", "earlier", "before"]),
                "refers_to_previous": any(kw in query.lower() for kw in ["previous", "earlier", "before"]),
                "needs_retrieval": any(kw in query.lower() for kw in ["document", "file", "pdf", "uploaded"])
            }

    def handle_general_query(self, query: str) -> str:
        history = self.memory.get_recent_history()
        task_description = (
            f"Jawab pertanyaan berikut dengan mempertimbangkan konteks percakapan:\n\n"
            f"Percakapan:\n{history}\n\n"
            f"Pertanyaan: {query}\n\n"
            f"Berikan jawaban komprehensif yang kontekstual. \n\n"
            "Gunakan format yang tepat tanpa Markdown atau simbol seperti '**' untuk bold."
        )
        response_task = Task(
            description=task_description,
            agent=self.agent,
            expected_output="Context-aware comprehensive answer"
        )
        try:
            return self.agent.execute_task(response_task)
        except Exception as e:
            logging.error(f"Error saat menangani kueri umum: {e}")
            return f"Error: {e}"
    
    def handle_rag_query(self, query: str) -> str:
        """Menangani query yang membutuhkan retrieval dari dokumen"""
        # Dapatkan konteks relevan dari RAG
        context = rag_manager.get_relevant_context(query)
        
        history = self.memory.get_recent_history(2)
        task_description = (
            f"Jawab pertanyaan berikut berdasarkan konteks dokumen yang disediakan:\n\n"
            f"Percakapan terakhir:\n{history}\n\n"
            f"Pertanyaan: {query}\n\n"
            f"Konteks dokumen:\n{context}\n\n"
            "Berikan jawaban komprehensif berdasarkan konteks dokumen yang diberikan. "
            "Jika jawabannya tidak ada dalam konteks, jawab jujur bahwa informasi tersebut "
            "tidak ditemukan dalam dokumen yang tersedia."
            "Gunakan format yang tepat tanpa Markdown atau simbol seperti '**' untuk bold."
        )
        
        rag_task = Task(
            description=task_description,
            agent=self.agent,
            expected_output="Document-based comprehensive answer"
        )
        
        try:
            return self.agent.execute_task(rag_task)
        except Exception as e:
            logging.error(f"Error saat melakukan RAG query: {e}")
            return f"Error saat mencari informasi dalam dokumen: {e}"

    def get_agent_response(self, query: str, agent_id: str, with_context: bool = True) -> str:
        if agent_id not in agents_db:
            error_msg = f"Agen '{agent_id}' tidak ditemukan."
            logging.error(error_msg)
            return error_msg
            
        # Tambahkan konteks RAG jika diperlukan
        context = ""
        if with_context:
            history = self.memory.get_recent_history()
            context = f"Percakapan:\n{history}\n\n"
        
        # Coba dapatkan konteks dokumen
        doc_context = rag_manager.get_relevant_context(query, k=3)
        if doc_context:
            context += f"Konteks dokumen relevan:\n{doc_context}\n\n"
        
        task_description = (
            f"Sebagai agen, jawab pertanyaan berikut:\n\n"
            f"{context}"
            f"Pertanyaan: {query}\n\n"
            "Berikan jawaban teknis dan mendalam, memanfaatkan konteks dokumen jika relevan."
            "Gunakan format yang tepat tanpa Markdown atau simbol seperti '**' untuk bold."
        )
        
        custom_task = Task(
            description=task_description,
            agent=agents_db[agent_id],
            expected_output="Expert domain answer"
        )
        
        try:
            return agents_db[agent_id].execute_task(custom_task)
        except Exception as e:
            logging.error(f"Error pada respons agen: {e}")
            return f"Error dari agen {agent_id}: {e}"

    def process_query(self, query: str, session_id: Optional[str] = None) -> str:
        if session_id and session_id != self.session_id:
            self.session_id = session_id
            self.memory = ConversationMemory(session_id)
            
        analysis = self.analyze_query(query)
        
        # Jika query membutuhkan retrieval dokumen
        if analysis.get("needs_retrieval", False) or "dokumen" in query.lower():
            response = self.handle_rag_query(query)
        # Replace with:
        elif analysis["query_type"] in ["general", "follow_up"]:
            response = self.handle_general_query(query)
        else:
            # Handle case where agent_needed might be a boolean
            role_type = analysis.get("agent_needed")
            if isinstance(role_type, bool):
                # Default to programming agent if agent_needed is just True
                role_type = "programming"
            elif not role_type:
                # Default for None or False
                role_type = "programming"
                
            agent_id = f"{role_type.lower().replace(' ', '_')}_agent"
            agent_response = self.get_agent_response(query, agent_id, with_context=analysis.get("requires_context", False))
            agent_responses = {agent_id: agent_response}
            response = self.aggregate_responses(query, agent_response, agent_responses)
            
        self.memory.add_interaction(query, response)
        if len(conversation_db[self.session_id]) % 5 == 0:
            self.memory.save_to_file()
        return response
        
    def aggregate_responses(self, query: str, main_response: str, agent_responses: dict) -> str:
        if not agent_responses:
            return main_response
        agents_input = "\n".join([f"- {aid}: {resp}" for aid, resp in agent_responses.items()])
        history = self.memory.get_recent_history(2)
        task_description = (
            f"Agregasikan respons berikut menjadi jawaban yang koheren untuk kueri:\n\n"
            f"Percakapan terakhir:\n{history}\n\n"
            f"Kueri: \"{query}\"\n\n"
            f"Respons agen:\n{agents_input}\n\n"
            f"Berikan jawaban yang tersintesis, kontekstual, dan komprehensif.\n\n"
            f"Gunakan format yang tepat tanpa Markdown atau simbol seperti '**' untuk bold."
        )
        aggregation_task = Task(
            description=task_description,
            agent=self.agent,
            expected_output="Context-aware synthesized comprehensive answer"
        )
        try:
            return self.agent.execute_task(aggregation_task)
        except Exception as e:
            logging.error(f"Error saat mengagregasi respons: {e}")
            return main_response


    def create_dynamic_task(self, domain, query_type, query, context=None, parameters=None):
        """Generate dynamic task descriptions based on query type and context"""
        base_prompts = {
            "general": "Jawab pertanyaan umum berikut dengan pengetahuan yang luas:",
            "follow_up": "Jawab pertanyaan lanjutan ini dengan mempertimbangkan konteks sebelumnya:",
            "needs_agent": f"Sebagai spesialis {domain}, jawab pertanyaan berikut dengan keahlian teknis:",
            "needs_retrieval": "Jawab pertanyaan ini berdasarkan dokumen yang diambil:"
        }
        
        # Start with base prompt for this query type
        task_prompt = base_prompts.get(query_type, base_prompts["general"])
        
        # Add context if available
        if context:
            task_prompt += f"\n\nKonteks:\n{context}"
        
        # Add query
        task_prompt += f"\n\nPertanyaan: {query}"
        
        # Add any additional parameters or instructions
        if parameters:
            for key, value in parameters.items():
                task_prompt += f"\n\n{key}: {value}"
        
        return task_prompt

    def create_specialized_agent(self, domain, personality=None, expertise_level="expert", model=None):
        """Create a specialized agent with custom traits"""
        
        # Default model if none provided
        if not model:
            model = self.model
        
        # Define personality traits
        personalities = {
            "teacher": "Anda menjelaskan dengan cara yang jelas dan sederhana, membuat konsep kompleks mudah dipahami.",
            "academic": "Anda memberikan jawaban komprehensif dengan referensi ilmiah dan nuansa akademis.",
            "practical": "Anda fokus pada solusi praktis dan langkah-langkah implementasi nyata.",
            "creative": "Anda berpikir out-of-the-box dan menawarkan perspektif inovatif."
        }
        
        # Use provided personality or default
        backstory = personalities.get(personality, "Anda adalah ahli teknis dengan pengetahuan mendalam.")
        backstory += f" Sebagai {expertise_level} dalam {domain}, Anda memberikan jawaban akurat dan terpercaya."
        
        agent_llm = LLM(model=model, api_key=GEMINI_API_KEY)
        
        custom_agent = Agent(
            role=f"{domain.capitalize()} Specialist",
            goal=f"Memberikan jawaban optimal tentang {domain} dengan tingkat keahlian {expertise_level}",
            backstory=backstory,
            verbose=True,
            allow_delegation=False,
            llm=agent_llm
        )
        
        # Generate unique ID
        agent_id = f"{domain.lower().replace(' ', '_')}_{personality or 'standard'}_{datetime.now().strftime('%Y%m%d%H%M%S')}"
        agents_db[agent_id] = custom_agent
        
        return agent_id, custom_agent

    def enhanced_query_analyzer(self, query, history=None):
        """Advanced query analyzer that detects specific domains, skills, and requirements"""
        
        # Include conversation history for context
        if history is None:
            history = self.memory.get_recent_history(3)
        
        analysis_prompt = (
            f"Analisis kueri berikut secara mendalam:\n\n"
            f"Riwayat percakapan:\n{history}\n\n"
            f"Kueri: \"{query}\"\n\n"
            "Berikan analisis komprehensif dengan JSON yang mencakup:\n"
            "1. query_type: ['general', 'follow_up', 'technical', 'creative', 'opinion', 'factual']\n"
            "2. domain_primary: domain utama (programming, science, math, business, dll)\n"
            "3. domain_secondary: sub-domain spesifik jika ada\n" 
            "4. complexity_level: [1-5] dengan 5 sebagai yang paling kompleks\n"
            "5. requires_context: [true/false] apakah membutuhkan konteks sebelumnya\n"
            "6. needs_retrieval: [true/false] apakah membutuhkan pencarian dokumen\n"
            "7. sentiment: sentimen pengguna ['positive', 'negative', 'neutral', 'urgent']\n"
            "8. requires_specialization: [true/false] apakah membutuhkan pengetahuan khusus\n"
            "9. expected_format: format jawaban yang diharapkan ('text', 'code', 'list', 'tutorial')\n"
            "Gunakan format JSON yang tepat tanpa komentar tambahan."
        )
        
        analysis_task = Task(
            description=analysis_prompt,
            agent=self.agent,
            expected_output="Detailed JSON analysis of the query"
        )
        
        try:
            result = self.agent.execute_task(analysis_task)
            json_start = result.find("{")
            json_end = result.rfind("}") + 1
            if json_start == -1 or json_end == -1:
                raise ValueError("Output JSON tidak ditemukan.")
            
            analysis = json.loads(result[json_start:json_end])
            return analysis
        except Exception as e:
            logging.error(f"Error saat analisis kueri: {e}")
            # Fallback analysis
            return {
                "query_type": "general",
                "domain_primary": "general",
                "domain_secondary": None, 
                "complexity_level": 3,
                "requires_context": any(kw in query.lower() for kw in ["sebelumnya", "tadi", "lanjutkan"]),
                "needs_retrieval": any(kw in query.lower() for kw in ["dokumen", "file", "pdf"]),
                "sentiment": "neutral",
                "requires_specialization": False,
                "expected_format": "text"
            }
        
    def process_dynamic_query(self, query, session_id=None):
        """Process queries with dynamic task creation and agent selection"""
        
        # Get or initialize session
        if session_id and session_id != self.session_id:
            self.session_id = session_id
            self.memory = ConversationMemory(session_id)
        
        # Analyze query
        analysis = self.enhanced_query_analyzer(query)
        
        # Determine if we need to create a new specialized agent
        if analysis.get("requires_specialization") and analysis.get("domain_primary") != "general":
            domain = analysis.get("domain_primary")
            domain_id = f"{domain}_agent"
            
            # Check if we already have an agent for this domain
            if domain_id not in agents_db:
                # Create specialized agent dynamically
                personality = "teacher" if analysis.get("complexity_level", 3) <= 2 else "academic"
                domain_id, _ = self.create_specialized_agent(domain, personality)
        
        # Handle retrieval if needed
        context = ""
        if analysis.get("needs_retrieval"):
            context = rag_manager.get_relevant_context(query)
        
        # Include conversation history if needed
        if analysis.get("requires_context"):
            history = self.memory.get_recent_history()
            context = f"{history}\n\n{context}" if context else history
        
        # Create dynamic task
        task_params = {
            "expected_format": analysis.get("expected_format", "text"),
            "complexity_level": analysis.get("complexity_level", 3)
        }
        
        task_description = self.create_dynamic_task(
            analysis.get("query_type", "general"),
            query,
            context=context,
            parameters=task_params
        )
        
        # Select appropriate agent
        if analysis.get("requires_specialization") and analysis.get("domain_primary") != "general":
            agent_id = f"{analysis.get('domain_primary')}_agent"
            agent = agents_db.get(agent_id, self.agent)
        else:
            agent = self.agent
        
        # Execute task
        dynamic_task = Task(
            description=task_description,
            agent=agent,
            expected_output=f"{analysis.get('expected_format')} response at complexity level {analysis.get('complexity_level')}"
        )
        
        try:
            response = agent.execute_task(dynamic_task)
        except Exception as e:
            logging.error(f"Error executing task: {e}")
            response = f"Maaf, terjadi kesalahan saat memproses permintaan Anda: {e}"
        
        # Save interaction
        self.memory.add_interaction(query, response, agent_id=getattr(agent, 'role', 'orchestrator'))
        
        return response
    
    def update_agent_knowledge(self, agent_id, feedback=None, new_information=None):
        """Update agent knowledge based on feedback or new information"""
        if agent_id not in agents_db:
            return False
        
        agent = agents_db[agent_id]
        
        # Store information in agent's knowledge base
        if not hasattr(agent, "knowledge_base"):
            agent.knowledge_base = []
        
        # Add feedback for learning
        if feedback:
            agent.knowledge_base.append({
                "type": "feedback",
                "content": feedback,
                "timestamp": datetime.now().isoformat()
            })
        
        # Add new information
        if new_information:
            agent.knowledge_base.append({
                "type": "information",
                "content": new_information,
                "timestamp": datetime.now().isoformat()
            })
        
        # Update agent's backstory to include new knowledge
        if hasattr(agent, "knowledge_base") and len(agent.knowledge_base) > 0:
            # Extract recent knowledge items
            recent_items = agent.knowledge_base[-5:]
            knowledge_summary = "\n".join([f"- {item['content']}" for item in recent_items])
            
            # Update backstory with new knowledge
            updated_backstory = f"{agent.backstory}\n\nRecent knowledge updates:\n{knowledge_summary}"
            agent.backstory = updated_backstory
        
        return True
            
def extract_text_from_pdf(pdf_path):
    """Extract text from a PDF file using PyPDF2."""
    try:
        with open(pdf_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            text = "\n".join([page.extract_text() for page in pdf_reader.pages if page.extract_text()])
        return text.strip() if text else "Error: PDF tidak mengandung teks yang dapat diekstrak."
    except Exception as e:
        return f"Error extracting PDF: {e}"

def extract_text_from_docx(docx_path):
    """Extract text from a DOCX file using python-docx."""
    try:
        doc = docx.Document(docx_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip() if text else "Error: DOCX kosong atau tidak dapat diekstrak."
    except Exception as e:
        return f"Error extracting DOCX: {e}"

def extract_text_from_txt(txt_path):
    """Extract text from a TXT file with encoding detection."""
    try:
        with open(txt_path, "rb") as f:
            raw_data = f.read()
            encoding = chardet.detect(raw_data)["encoding"] or "utf-8"

        with open(txt_path, "r", encoding=encoding, errors="ignore") as f:
            return f.read().strip()
    except Exception as e:
        return f"Error extracting TXT: {e}"

def create_custom_agent(domain: str, model: str = "gemini/gemini-1.5-flash") -> str:
    my_llm = LLM(model=model, api_key=GEMINI_API_KEY)
    custom_agent = Agent(
        role=f"{domain.capitalize()} Agent",
        goal=f"Memberikan pengetahuan mendalam tentang {domain}",
        backstory=f"Anda adalah ahli dalam {domain} dengan pengetahuan teknis mendalam.",
        verbose=True,
        allow_delegation=False,
        llm=my_llm
    )
    agent_id = f"{domain.lower().replace(' ', '_')}_agent"
    agents_db[agent_id] = custom_agent
    logging.info(f"Agen baru dibuat: {agent_id}")
    return agent_id

def save_agents_to_file():
    agents_data = {
        name: {
            "role": agent.role,
            "goal": agent.goal,
            "backstory": agent.backstory,
            "model": agent.llm.model
        }
        for name, agent in agents_db.items()
    }
    with open("agents_db.json", "w") as f:
        json.dump(agents_data, f, indent=2)
    logging.info("Agen berhasil disimpan ke agents_db.json")

def load_agents_from_file():
    if os.path.exists("agents_db.json"):
        with open("agents_db.json", "r") as f:
            data = json.load(f)
            for name, details in data.items():
                my_llm = LLM(model=details["model"], api_key=GEMINI_API_KEY)
                agents_db[name] = Agent(
                    role=details["role"],
                    goal=details["goal"],
                    backstory=details["backstory"],
                    llm=my_llm,
                    verbose=True,
                    allow_delegation=("orchestrator" in name)
                )
        logging.info(f"Berhasil memuat {len(agents_db)} agen dari file")


def get_or_create_session(session_id: str = "default") -> OrchestratorAgent:
    if session_id in sessions:
        return sessions[session_id]
    if "orchestrator" not in agents_db:
        initialize_system()
    new_orchestrator = OrchestratorAgent(session_id=session_id)
    if "orchestrator" in agents_db:
        new_orchestrator.agent = agents_db["orchestrator"]
    sessions[session_id] = new_orchestrator
    return new_orchestrator

def process_query(query: str, session_id: str = "default", file_path: str = None) -> dict:
    """
    Process user query and return response with additional metadata.
    
    Returns a dictionary containing:
    - response: text response
    - downloadableFile: (optional) information about a file to download
    """
    session_orchestrator = get_or_create_session(session_id)
    
    try:
        document_text = ""
        if file_path:
            # Extract text from uploaded document
            if file_path.endswith(".pdf"):
                document_text = extract_text_from_pdf(file_path)
            elif file_path.endswith(".docx"):
                document_text = extract_text_from_docx(file_path)
            elif file_path.endswith(".txt"):
                document_text = extract_text_from_txt(file_path)
            else:
                return {"response": "Error: Format file tidak didukung."}
        
        # Append document content to query if available
        if document_text:
            query = f"{query}\n\n[Dokumen Referensi]\n{document_text[:1000]}..."
        
        # Check if the query is for document conversion
        convert_keywords = ["convert to word", "convert", "konversi", "konversi ke word", "buatlah proposal"]
        
        # Ensure query is a string before processing
        if not isinstance(query, str):
            return {"response": "Error: Query tidak valid."}
        
        # Handle document conversion request
        is_conversion_request = any(keyword in query.lower() for keyword in convert_keywords)
        
        if is_conversion_request:
            # Retrieve conversation history
            history = conversation_db.get(session_id, [])
            
            if len(history) < 2:
                return {"response": "Tidak ada jawaban sebelumnya untuk dikonversi."}
            
            # Get the latest agent response (not the current one being processed)
            latest_response = None
            for entry in reversed(history):
                if isinstance(entry, dict) and "agent_response" in entry:
                    latest_response = entry["agent_response"]
                    break
            
            if not latest_response:
                return {"response": "Tidak ada konten yang dapat dikonversi."}
            
            # Determine file format based on query
            file_format = "docx" if "word" in query.lower() else "pdf"
            filename = f"response_{session_id}.{file_format}"
            filepath = os.path.join("outputs", filename)
            
            # Create directory if it doesn't exist
            os.makedirs("outputs", exist_ok=True)
            
            # Save the content to a file
            if file_format == "docx":
                save_to_word(latest_response, filepath)
            else:
                save_to_pdf(latest_response, filepath)
            
            # Return both response and downloadable file info
            return {
                "response": f"Dokumen telah dibuat dan siap diunduh.",
                "downloadableFile": {
                    "url": f"/download/{filename}",
                    "filename": filename
                }
            }
        
        # Normal query processing
        response = session_orchestrator.process_query(query, session_id)
        return {"response": response}
    
    except Exception as e:
        logging.error(f"Error saat memproses kueri: {e}")
        return {"response": f"Error saat memproses kueri: {e}"}

def create_custom_agent(domain: str, model: str = "gemini/gemini-1.5-flash") -> str:
    my_llm = LLM(model=model, api_key=GEMINI_API_KEY)
    custom_agent = Agent(
        role=f"{domain.capitalize()} Agent",
        goal=f"Memberikan pengetahuan mendalam tentang {domain}",
        backstory=f"Anda adalah ahli dalam {domain} dengan pengetahuan teknis mendalam.",
        verbose=True,
        allow_delegation=False,
        llm=my_llm
    )
    agent_id = f"{domain.lower().replace(' ', '_')}_agent"
    agents_db[agent_id] = custom_agent
    logging.info(f"Agen baru dibuat: {agent_id}")
    return agent_id

def initialize_system() -> bool:
    logging.info("Inisialisasi sistem agen...")
    try:
        if "orchestrator" not in agents_db:
            logging.info("Membuat orchestrator agent...")
            OrchestratorAgent()
            logging.info("Orchestrator agent berhasil dibuat.")
        domain_types = {
            "programming": "Memberikan panduan pengembangan perangkat lunak dan pemrograman",
            "science": "Menjelaskan konsep ilmiah dengan akurasi dan kejelasan"
        }
        for domain, goal in domain_types.items():
            agent_id = f"{domain}_agent"
            if agent_id not in agents_db:
                logging.info(f"Membuat {domain} agent...")
                my_llm = LLM(model="gemini/gemini-1.5-flash", api_key=GEMINI_API_KEY)
                domain_agent = Agent(
                    role=f"{domain.capitalize()} Agent",
                    goal=goal,
                    backstory=f"Anda adalah ahli dalam {domain} dengan pengetahuan mendalam.",
                    verbose=True,
                    allow_delegation=False,
                    llm=my_llm
                )
                agents_db[agent_id] = domain_agent
                logging.info(f"{domain.capitalize()} agent berhasil dibuat.")
        save_agents_to_file()
        logging.info(f"Sistem diinisialisasi dengan {len(agents_db)} agen.")
        if "orchestrator" not in agents_db:
            logging.error("Gagal membuat orchestrator!")
            return False
        return True
    except Exception as e:
        logging.error(f"ERROR saat inisialisasi: {e}")
        try:
            OrchestratorAgent()
            logging.info("Fallback: Orchestrator berhasil dibuat setelah error.")
            return True
        except Exception as ex:
            logging.error(f"Gagal membuat fallback orchestrator: {ex}")
            return False

def process_document(file_path, file_type, metadata=None):
    if metadata is None:
        metadata = {"source": "uploaded_document"}
    
    # Pilih loader sesuai file_type
    if file_type == "pdf":
        loader = PyPDFLoader(file_path)
    elif file_type in ["doc", "docx"]:
        loader = UnstructuredWordDocumentLoader(file_path)
    elif file_type in ["txt", "md"]:
        loader = TextLoader(file_path)
    else:
        return "Unsupported file format"
    
    documents = loader.load()
    text_splitter = CharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    docs = text_splitter.split_documents(documents)
    
    # Perbarui metadata di setiap dokumen
    for doc in docs:
        doc.metadata.update(metadata)

    vectorstore.add_documents(docs)

    return "Knowledge added successfully"

# ------------------------------------------
# Implementasi Flask
# ------------------------------------------
app = Flask(__name__)

app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

app.secret_key = os.environ.get("SECRET_KEY")     

CORS(app, supports_credentials=True, resources={
    r"/*": {
        "origins": ["*"],
    }
})

@app.route('/')
def home():
    return """
    <h1>Agent System</h1>
    <p>Gunakan endpoint /chat untuk mengirim kueri.</p>
    """

@app.route('/chat', methods=['POST'])
def chat():
    query = request.form.get("query", "").strip()  # Get query from form
    session_id = request.form.get("session_id", "default")

    if not query and 'document' not in request.files:
        return jsonify({"error": "Query kosong atau tidak ada dokumen."}), 400

    file_path = None
    if 'document' in request.files:
        file = request.files['document']
        if file.filename:
            filename = secure_filename(file.filename)
            file_path = os.path.join(UPLOAD_FOLDER, filename)
            file.save(file_path)
            
            try:
                extracted_text = extract_text_from_document(file_path)
                
                metadata = {
                    "source": file.filename,
                    "upload_time": datetime.now().isoformat(),
                    "file_type": os.path.splitext(file.filename)[1][1:]
                }
                
                rag_manager.add_document(extracted_text, metadata)
            
            except Exception as e:
                logging.error(f"Error saat mengunggah dan memproses dokumen: {e}")

    # Send query and document path to process_query
    response = process_query(query, session_id, file_path)

    # Clean up temporary file if it exists
    if file_path and os.path.exists(file_path):
        os.remove(file_path)
        
    # Handle both string and dictionary responses
    if isinstance(response, dict):
        return jsonify(response)
    else:
        # For backward compatibility with existing frontend
        return jsonify({"response": response})

@app.route("/download/<filename>", methods=['GET'])
def download_file(filename):
    file_path = os.path.join(OUTPUT_FOLDER, filename)
    app.logger.info(f"Attempting to download: {file_path}")
    app.logger.info(f"File exists: {os.path.exists(file_path)}")
    
    if os.path.exists(file_path):
        return send_file(file_path, as_attachment=True)
    else:
        return jsonify({"error": f"File {filename} not found"}), 404

# Create and Get Agent
@app.route('/agents', methods=['GET', 'POST'])
def agents():
    if request.method == 'GET':
        agents_list = [[agent_id, agent.role, agent.llm.model, agent.goal, agent.backstory,] for agent_id, agent in agents_db.items()]
        return jsonify({"agents": agents_list})
    
    elif request.method == 'POST':
        data = request.get_json()
        domain = data.get("domain", "")
        model = data.get("model", "gemini/gemini-1.5-flash")
        if not domain.strip():
            return jsonify({"error": "Domain agen harus diisi."}), 400
        agent_id = create_custom_agent(domain, model)
        save_agents_to_file()
        return jsonify({"message": f"Agen {domain} berhasil dibuat dengan ID: {agent_id}"}), 201

# Delete Agent
@app.route('/agents/<agent_id>', methods=['DELETE'])
def delete_agent(agent_id):
    if agent_id == "orchestrator":
        return jsonify({"error": "Tidak dapat menghapus agen orchestrator."}), 400
    if agent_id in agents_db:
        del agents_db[agent_id]
        save_agents_to_file()
        return jsonify({"message": f"Agen {agent_id} berhasil dihapus."})
    else:
        return jsonify({"error": f"Agen {agent_id} tidak ditemukan."}), 404

# Update Agent
@app.route('/agents/<agent_id>', methods=['PUT'])
def update_agent(agent_id):
    if agent_id not in agents_db:
        return jsonify({"error": "Agent not found"}), 404

    data = request.get_json()

    agent = agents_db[agent_id]
    agent.goal = data.get("goal", agent.goal)
    agent.backstory = data.get("backstory", agent.backstory)
    agent.llm.model = data.get("model", agent.llm.model)

    save_agents_to_file()

    return jsonify({"message": f"Agent {agent_id} updated successfully"}), 200

@app.route("/agents/<agent_id>/knowledge", methods=["POST"])
def upload_knowledge(agent_id):
    if "file" not in request.files:
        return jsonify({"error": "No file uploaded."}), 400
    
    file = request.files["file"]
    filename = secure_filename(file.filename)
    file_ext = filename.split(".")[-1].lower()
    
    if file_ext not in ["pdf", "doc", "docx", "txt", "md"]:
        return jsonify({"error": "Unsupported file format."}), 400
    
    save_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(save_path)
    result = process_document(save_path, file_ext)
    
    return jsonify({"message": result})

if __name__ == "__main__":
    logging.info("Memulai Agent Orchestrator System dengan Flask...")
    os.makedirs("conversations", exist_ok=True)
    try:
        load_agents_from_file()
    except Exception as e:
        logging.error(f"Error saat memuat agen dari file: {e}")
    if "orchestrator" not in agents_db:
        logging.info("Orchestrator tidak ditemukan, melakukan inisialisasi sistem...")
        if not initialize_system():
            logging.warning("Inisialisasi sistem mungkin gagal")
    if "orchestrator" in agents_db:
        orchestrator = OrchestratorAgent()
        orchestrator.agent = agents_db["orchestrator"]
        sessions["default"] = orchestrator
    logging.info(f"Agen aktif: {list(agents_db.keys())}")
    app.run(host="0.0.0.0", port=5000, debug=True)